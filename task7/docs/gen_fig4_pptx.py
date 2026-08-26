#!/usr/bin/env python3
"""把「三款四足機器狗比較」論文風圖表輸出成可編輯的 PPTX。

版面座標與 gen_fig4.py 完全一致（同一套 SVG 座標系），差別只在輸出後端：
每個文字都是獨立的文字方塊、每個框都是獨立的矩形，可在 PowerPoint／
LibreOffice Impress 裡直接點進去改字。

SVG px → PPT pt 的換算為 PT = 0.625（1840 px 寬 → 15.97 in 投影片）。
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

W = 1840
SERIF = "Noto Serif CJK TC"
MONO = "Noto Sans Mono CJK TC"

INK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x6B, 0x6B, 0x6B)
LINE = RGBColor(0x2B, 0x2B, 0x2B)
FILL_L = RGBColor(0xF2, 0xF1, 0xEE)
FILL_M = RGBColor(0xDC, 0xDA, 0xD5)
FILL_D = RGBColor(0xBC, 0xB9, 0xB1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

PT = 0.625                      # 1 SVG px = 0.625 pt


def P(px):
    """SVG px → pptx 長度。"""
    return Pt(px * PT)


prs = Presentation()
slide = None                    # 稍後決定尺寸後才建立


# ------------------------------------------------------------------ 繪圖基元
def est_w(s, size):
    """粗估字串寬度（SVG px）：中日韓全形算一個字寬，半形算 0.55。"""
    w = 0.0
    for ch in s:
        w += size if ord(ch) > 0x2000 else size * 0.55
    return w


def text(x, y, s, size=15, font=SERIF, weight="normal", anchor="start",
         fill=INK):
    """在 SVG 基線 (x, y) 放一個文字方塊；anchor: start / end。"""
    boxw = est_w(s, size) * 1.25 + 40            # 留餘裕，避免被自動折行
    left = x if anchor == "start" else x - boxw
    top = y - size * 1.15
    tb = slide.shapes.add_textbox(P(left), P(top), P(boxw), P(size * 1.9))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT if anchor == "start" else PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = s
    f = r.font
    f.name = font
    f.size = Pt(round(size * PT, 1))
    f.bold = (weight == "bold")
    f.color.rgb = fill
    # 中文字型要同時寫進 ea（East Asian）才不會被換掉
    rPr = r._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.makeelement(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            + tag.split(":")[1], {"typeface": font})
        rPr.append(el)
    return tb


def box(x, y, w, h, fill=WHITE, stroke=LINE, sw=1.2, dash=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, P(x), P(y), P(w), P(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if sw == 0 or stroke is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = stroke
        sh.line.width = Pt(max(0.5, sw * PT))
        if dash:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            sh.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    sh.text_frame.word_wrap = False
    return sh


def line(x1, y1, x2, y2, sw=1.0, stroke=LINE):
    """水平線用細矩形畫 —— 連接線在 Impress 會被套主題陰影。"""
    th = max(sw, 1.0)
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                P(x1), P(y1 - th / 2), P(x2 - x1), P(th))
    sh.shadow.inherit = False
    sh.fill.solid()
    sh.fill.fore_color.rgb = stroke
    sh.line.fill.background()
    return sh


# ------------------------------------------------------------------ 版面
L = 46
GUT = 128
COLW, GAP = 508, 26
CX = [L + GUT + i * (COLW + GAP) for i in range(3)]

HY = 90
HH = 76
heads = [
    ("D1 EDU", "輪足・小狗", "16 軸（12 腿＋4 輪）　20.6 kg"),
    ("D1 Max", "輪足・中狗", "16 軸（12 腿＋4 輪）　41 kg"),
    ("D1 MaxPro", "點足・大狗", "12 軸　68 kg"),
]

plats = [
    [("Firefly AIO-3588SJD4", "Rockchip RK3588・單板",
      ["CPU　4×A76 ＋ 4×A55", "RAM　7.7 GiB（無 swap）",
       "GPU　Mali-G610 MP4", "NPU　RKNPU driver 0.9.2",
       "OS　Ubuntu 22.04.4・eMMC 64 GB"], True)],
    [("Firefly AIO-3588SJD4", "Rockchip RK3588・運控",
      ["CPU　4×A76 ＋ 4×A55", "RAM　7.7 GiB",
       "GPU　Mali-G610 MP4", "NPU　RKNPU 0.9.2"], True),
     ("NVIDIA Jetson Orin NX 16GB", "應用・建圖／定位／導航",
      ["CPU　8×Cortex-A78AE　1.984 GHz", "RAM　15 GiB",
       "GPU　NVIDIA ga10b", "DLA　NVDLA0／NVDLA1"], True)],
    [("Firefly AIO-3588SJD4", "Rockchip RK3588・主控板",
      ["CPU　4×A76 ＋ 4×A55", "RAM　7.7 GiB（Swap 0）",
       "GPU　Mali-G610 MP4", "NPU　RKNPU 0.9.2・librknnrt／librga",
       "OS　Ubuntu 22.04.5・kernel 5.10-rt"], True),
     ("轉發板", "CAN／關節訊號轉發",
      ["非第二塊高算力板"], False)],
]


def board_h(rows):
    return 30 + 24 + len(rows) * 24 + 14


SY = HY + HH + 22
col_h = [sum(board_h(r) + 10 for (_, _, r, _) in bs) for bs in plats]
SH = max(col_h) + 38

levels = [
    "L4　單顆馬達　角度・角速度・力矩・Kp・Kd",
    "L3　關節層　位置／速度",
    "L2　姿態動作　站立・趴下・匍匐・爬階",
    "L1　整機速度　move(vx, vy, yaw)",
]
reach = [2, 2, 4]
LH = 40

TY = SY + SH + 30
ry = TY + 6 + 4 * LH + 10
VY = ry + 74
CH = 118
cy = VY + CH + 46
H = cy + 96

# ------------------------------------------------------------------ 建立投影片
prs.slide_width = P(W)
prs.slide_height = P(H)
slide = prs.slides.add_slide(prs.slide_layouts[6])   # 空白版面

# ------------------------------------------------------------------ 標題
text(L, 42, "智元三款四足機器狗：運算平台、SDK 能力與單顆馬達控制可行性", 23, weight="bold")
text(L, 68, "硬體規格為實機實測（cat /proc/device-tree/model、lscpu、free -h、dmesg）；"
            "SDK 與控制介面為官方文件＋二進位符號分析＋實機偵察", 13.5, fill=GREY)

# ------------------------------------------------------------------ 欄頭
for i, (nm, kind, spec) in enumerate(heads):
    x = CX[i]
    box(x, HY, COLW, HH, fill=FILL_M, sw=1.4)
    text(x + 16, HY + 30, nm, 21, weight="bold")
    text(x + 16 + len(nm) * 12.5 + 24, HY + 30, kind, 15, fill=GREY)
    text(x + 16, HY + 57, spec, 15)

# ------------------------------------------------------------------ 甲：運算平台
line(L, SY - 8, W - L, SY - 8, sw=1.4)
text(L, SY + 22, "（甲）", 16.5, weight="bold")
text(L, SY + 48, "運算平台", 17.5, weight="bold")
text(L, SY + 72, "實機實測", 12.5, fill=GREY)

for i, boards in enumerate(plats):
    x = CX[i]
    y = SY + 6
    for (nm, sub, rows, big) in boards:
        bh = board_h(rows)
        box(x, y, COLW, bh, fill=FILL_L if big else WHITE,
            sw=1.4 if big else 1.0, dash=None if big else "dash")
        mono = ("AIO" in nm) or ("Jetson" in nm)
        text(x + 16, y + 28, nm, 16.5, font=MONO if mono else SERIF, weight="bold")
        text(x + 16, y + 50, sub, 13.5, fill=GREY)
        for k, r in enumerate(rows):
            text(x + 16, y + 76 + k * 24, r, 14.5)
        y += bh + 10
    verdict = ["單一算力平台", "★ 兩塊獨立高算力板", "主控板＋轉發板（非第二算力板）"][i]
    text(x + 16, SY + SH + 6, verdict, 15,
         weight="bold" if i == 1 else "normal",
         fill=INK if i == 1 else GREY)

# ------------------------------------------------------------------ 乙：SDK 可寫深度
line(L, TY - 10, W - L, TY - 10, sw=1.4)
text(L, TY + 22, "（乙）", 16.5, weight="bold")
text(L, TY + 48, "SDK", 17.5, weight="bold")
text(L, TY + 72, "可寫深度", 17.5, weight="bold")
text(L, TY + 96, "官方介面", 12.5, fill=GREY)

for i in range(3):
    x = CX[i]
    for k, lab in enumerate(levels):
        yy = TY + 6 + k * LH
        depth = 4 - k
        ok = depth <= reach[i]
        box(x, yy, COLW, LH - 7, fill=FILL_D if ok else WHITE,
            sw=1.2 if ok else 0.9, dash=None if ok else "dash")
        text(x + 14, yy + 22, lab, 14.5, fill=INK if ok else GREY)
    stop_y = TY + 6 + (4 - reach[i]) * LH
    if reach[i] < 4:
        line(x, stop_y, x + COLW, stop_y, sw=2.6)
        lbl = "官方 SDK 止步於此"
        lw = len(lbl) * 13 + 12
        box(x + COLW - 8 - lw, stop_y - 25, lw, 21, fill=WHITE, sw=0)
        text(x + COLW - 12, stop_y - 9, lbl, 13.5, anchor="end", weight="bold")

for i in range(3):
    x = CX[i]
    rd = ["讀取　16 軸狀態・IMU",
          "讀取　16 軸狀態＋關節溫度・IMU・光達・導航",
          "讀取　12 軸狀態・IMU・里程計・點雲"][i]
    box(x, ry, COLW, 32, fill=WHITE, sw=0.9)
    text(x + 14, ry + 22, rd, 14, fill=GREY)

sdkname = ["mc_sdk::zsl_1w::HighLevel　UDP 50 Hz",
           "robot_sdk::SDKClient　UDP :8082",
           "高層 TCP ＋ 底層 ROS 1 topic"]
for i in range(3):
    text(CX[i] + 14, ry + 58, sdkname[i], 13, font=MONO, fill=GREY)

# ------------------------------------------------------------------ 丙：單顆馬達控制
line(L, VY - 10, W - L, VY - 10, sw=1.4)
text(L, VY + 22, "（丙）", 16.5, weight="bold")
text(L, VY + 48, "單顆馬達", 17.5, weight="bold")
text(L, VY + 72, "控制", 17.5, weight="bold")

cards = [
    ("官方不提供", "做得到", "/spline_shm　共享記憶體", ""),
    ("官方不提供", "極可能做得到", "/dev/shm/joint_cmd　16 軸 × 5 欄位", ""),
    ("★ 官方提供", "做得到", "ROS 1 topic　rt/lowcmd", "底層 SDK 須安裝 ROS 1"),
]
for i, (off, verd, path, note) in enumerate(cards):
    x = CX[i]
    strong = (i == 2)
    box(x, VY, COLW, CH, fill=FILL_L if strong else WHITE, sw=1.6 if strong else 1.2)
    text(x + 16, VY + 28, off, 15.5, weight="bold" if strong else "normal",
         fill=INK if strong else GREY)
    text(x + 16, VY + 60, verd, 21, weight="bold")
    text(x + 16, VY + 88, path, 14.5, font=MONO)
    if note:
        text(x + COLW - 16, VY + 28, note, 13, anchor="end", fill=GREY)

# ------------------------------------------------------------------ 圖說
line(L, cy - 20, W - L, cy - 20, sw=1.4)
text(L, cy, "圖 4 · 三款機型比較。", 14.5, weight="bold")
text(L + 152, cy,
     "（甲）運算平台：D1 Max 是三台中唯一的雙高算力板；D1 MaxPro 雖是大狗，"
     "算力平台與小狗 D1 EDU 同為單顆 RK3588。", 14.5)
text(L, cy + 23,
     "（乙）三台的「讀」幾乎全開，差別在「寫」——D1 EDU 與 D1 Max 的官方 SDK 只到姿態動作層，"
     "送不了關節角度或力矩；只有 D1 MaxPro 開放到單顆馬達。", 14.5)
text(L, cy + 46,
     "（丙）D1 EDU 與 D1 Max 官方雖不提供，但機上存在共享記憶體介面可繞過 SDK："
     "D1 EDU 已實機驗證三種控制模式，D1 Max 結構已解出、僅差寫入未測。", 14.5)
text(L, cy + 72,
     "註：D1 MaxPro 官方文件的 SSH 帳號為 jetson，但實測主控板為 Firefly RK3588，"
     "兩者不一致，機型批次差異待向原廠確認。", 12.5, fill=GREY)

out = Path(__file__).resolve().parent / "圖4_三機型比較_論文風.pptx"
prs.save(out)
print("PPTX", out, f"{out.stat().st_size} bytes",
      f"slide {prs.slide_width.inches:.2f} × {prs.slide_height.inches:.2f} in")
