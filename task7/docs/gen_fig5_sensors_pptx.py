#!/usr/bin/env python3
"""把「三款四足機器狗：感測器可讀資料與格式」輸出成可編輯的一頁 PPTX。

版面座標與 gen_fig5_sensors.py 完全一致（同一套 SVG 座標系），差別只在輸出後端：
每個文字都是獨立的文字方塊、每個框都是獨立的矩形，可在 PowerPoint／
LibreOffice Impress 裡直接點進去改字。

SVG px → PPT pt 的換算為 PT = 0.625（1840 px 寬 → 15.97 in 投影片）。
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

W = 1840
SERIF = "Noto Serif CJK TC"
MONO = "Noto Sans Mono CJK TC"

INK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x6B, 0x6B, 0x6B)
LINE = RGBColor(0x2B, 0x2B, 0x2B)
FILL_L = RGBColor(0xF2, 0xF1, 0xEE)
FILL_M = RGBColor(0xDC, 0xDA, 0xD5)
FILL_D = RGBColor(0xBC, 0xB9, 0xB1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

PT = 0.625                      # 1 SVG px = 0.625 pt


def P(px):
    """SVG px → pptx 長度。"""
    return Pt(px * PT)


prs = Presentation()
slide = None                    # 稍後決定尺寸後才建立


# ------------------------------------------------------------------ 繪圖基元
def est_w(s, size):
    """粗估字串寬度（SVG px）：中日韓全形算一個字寬，半形算 0.55。"""
    w = 0.0
    for ch in s:
        w += size if ord(ch) > 0x2000 else size * 0.55
    return w


def text(x, y, s, size=15, font=SERIF, weight="normal", anchor="start",
         fill=INK, ls=None):
    """在 SVG 基線 (x, y) 放一個文字方塊；anchor: start / end。"""
    if slide is None:
        return None
    boxw = est_w(s, size) * 1.25 + 40            # 留餘裕，避免被自動折行
    left = x if anchor == "start" else x - boxw
    top = y - size * 1.15
    tb = slide.shapes.add_textbox(P(left), P(top), P(boxw), P(size * 1.9))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT if anchor == "start" else PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = s
    f = r.font
    f.name = font
    f.size = Pt(round(size * PT, 1))
    f.bold = (weight == "bold")
    f.color.rgb = fill
    # 中文字型要同時寫進 ea（East Asian）才不會被換掉
    rPr = r._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.makeelement(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            + tag.split(":")[1], {"typeface": font})
        rPr.append(el)
    return tb


def box(x, y, w, h, fill=WHITE, stroke=LINE, sw=1.2, dash=None, rx=3):
    if slide is None:
        return None
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, P(x), P(y), P(w), P(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if sw == 0 or stroke is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = stroke
        sh.line.width = Pt(max(0.5, sw * PT))
        if dash:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            sh.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    sh.text_frame.word_wrap = False
    return sh


def line(x1, y1, x2, y2, sw=1.0, stroke=LINE, dash=None):
    """水平線用細矩形畫 —— 連接線在 Impress 會被套主題陰影。"""
    if slide is None:
        return None
    th = max(sw, 1.0)
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                P(x1), P(y1 - th / 2), P(x2 - x1), P(th))
    sh.shadow.inherit = False
    sh.fill.solid()
    sh.fill.fore_color.rgb = stroke
    sh.line.fill.background()
    return sh




# ------------------------------------------------------------------ 版面
# ★ 版面**不在這裡重寫一份** —— 直接抓 gen_fig5_sensors.py 的版面段來跑。
#   之前這裡是拼接出來的複本，改了 SVG 卻忘了重拼，就會靜默產生過期的 PPTX
#   （2026-08-27 實際發生過一次）。現在兩邊永遠是同一份程式。
SRC = Path(__file__).resolve().parent / "gen_fig5_sensors.py"
_s = SRC.read_text(encoding="utf-8")
_marker = "# ------------------------------------------------------------------ 版面"
if _marker not in _s or "H = yy + 34" not in _s:
    raise SystemExit(f"❌ {SRC.name} 的版面段標記找不到 —— 兩支已經對不上，先修再跑")
LAYOUT = compile(_s.split(_marker, 1)[1].split("H = yy + 34")[0]
                 .replace('A("@@SVGOPEN@@")\n', ''),
                 str(SRC), "exec")


def draw():
    """跑一趟版面；slide is None 時所有基元都是 no-op，只用來先量高度。"""
    ns = dict(text=text, box=box, line=line, W=W, SERIF=SERIF, MONO=MONO,
              INK=INK, GREY=GREY, LINE=LINE, FILL_L=FILL_L, FILL_M=FILL_M,
              FILL_D=FILL_D, WHITE=WHITE)
    exec(LAYOUT, ns)
    return ns["yy"] + 34


H = draw()                                               # 第一趟：只量高度
prs.slide_width = P(W)
prs.slide_height = P(H)
slide = prs.slides.add_slide(prs.slide_layouts[6])       # 空白版面
draw()                                                   # 第二趟：真的畫

out = Path(__file__).resolve().parent / "圖5_三機型感測器_論文風.pptx"
prs.save(out)
print("PPTX", out, f"{out.stat().st_size} bytes",
      f"slide {prs.slide_width.inches:.2f} × {prs.slide_height.inches:.2f} in")
